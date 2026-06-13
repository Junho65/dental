from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_PATH = Path("artifacts/presentations/Dental_AI_Project_submission_report.pptx")


BG = RGBColor(246, 248, 251)
NAVY = RGBColor(20, 38, 66)
TEAL = RGBColor(26, 125, 140)
ORANGE = RGBColor(232, 115, 69)
TEXT = RGBColor(34, 43, 56)
MUTED = RGBColor(96, 109, 128)
LINE = RGBColor(211, 219, 229)
WHITE = RGBColor(255, 255, 255)


def set_background(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, section: str) -> None:
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.42),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = section
    p.font.name = "Aptos"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.32)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.72), Inches(11.8), Inches(0.9))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = "Malgun Gothic"
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED
        p2.space_before = Pt(6)


def add_bullet_box(slide, left, top, width, height, title: str, bullets: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.1)

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEAL

    for bullet in bullets:
        bp = tf.add_paragraph()
        bp.text = bullet
        bp.font.name = "Malgun Gothic"
        bp.font.size = Pt(11.5)
        bp.font.color.rgb = TEXT
        bp.level = 0
        bp.bullet = True
        bp.space_before = Pt(3)


def add_paragraph_box(slide, left, top, width, height, title: str, paragraphs: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.1)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.08)

    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Malgun Gothic"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = TEAL

    for text in paragraphs:
        pp = tf.add_paragraph()
        pp.text = text
        pp.font.name = "Malgun Gothic"
        pp.font.size = Pt(11.2)
        pp.font.color.rgb = TEXT
        pp.space_before = Pt(5)


def add_pipeline_box(slide, top, labels: list[str]) -> None:
    left = Inches(0.8)
    box_w = Inches(2.2)
    gap = Inches(0.22)
    arrow_w = Inches(0.35)

    for i, label in enumerate(labels):
        x = left + i * (box_w + arrow_w + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, box_w, Inches(0.95))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = TEAL if i % 2 == 0 else ORANGE
        shape.line.width = Pt(2)
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = "Malgun Gothic"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = TEXT
        p.alignment = PP_ALIGN.CENTER
        if i < len(labels) - 1:
            arrow = slide.shapes.add_textbox(x + box_w + Inches(0.03), top + Inches(0.25), arrow_w, Inches(0.35))
            atf = arrow.text_frame
            atf.clear()
            ap = atf.paragraphs[0]
            ap.text = "→"
            ap.font.name = "Aptos"
            ap.font.size = Pt(22)
            ap.font.bold = True
            ap.font.color.rgb = MUTED
            ap.alignment = PP_ALIGN.CENTER


def add_table_slide(slide, left, top, width, height, rows, cols, data, col_widths=None) -> None:
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    if col_widths:
        for idx, w in enumerate(col_widths):
            table.columns[idx].width = w
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(233, 241, 245) if r == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = "Malgun Gothic"
                paragraph.font.size = Pt(10.3 if r else 10.8)
                paragraph.font.bold = r == 0
                paragraph.font.color.rgb = NAVY if r == 0 else TEXT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_footer(slide, page_num: int) -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.7),
        Inches(7.0),
        Inches(11.9),
        Inches(0.01),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(11.95), Inches(7.05), Inches(0.5), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.font.name = "Aptos"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1. Title
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    add_title(slide, "치과 X-ray 병변 자동 검출 및 진단 보조 시스템", "제출용 프로젝트 정리 문서")
    title_box = slide.shapes.add_textbox(Inches(0.78), Inches(1.65), Inches(11.3), Inches(2.0))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Dental AI Project"
    p.font.name = "Aptos"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p2 = tf.add_paragraph()
    p2.text = "파노라마 치과 X-ray에서 병변 위치를 자동 검출하고,\n충치 계열 응답 정규화와 치주 후속 분류를 결합한 설명 가능한 진단 보조를 제공한다."
    p2.font.name = "Malgun Gothic"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.space_before = Pt(10)
    info = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.8), Inches(5.5), Inches(1.15))
    info.fill.solid()
    info.fill.fore_color.rgb = WHITE
    info.line.color.rgb = LINE
    itf = info.text_frame
    itf.clear()
    ip = itf.paragraphs[0]
    ip.text = "학번 21012031 | 문준호"
    ip.font.name = "Malgun Gothic"
    ip.font.size = Pt(18)
    ip.font.bold = True
    ip.font.color.rgb = TEXT
    ip.alignment = PP_ALIGN.CENTER
    ip2 = itf.add_paragraph()
    ip2.text = "딥러닝실습 프로젝트 제출 자료"
    ip2.font.name = "Malgun Gothic"
    ip2.font.size = Pt(12)
    ip2.font.color.rgb = MUTED
    ip2.alignment = PP_ALIGN.CENTER

    # 2. Overview
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "OVERVIEW")
    add_title(slide, "프로젝트 개요", "문제의식, 목표, 사용자 가치")
    add_bullet_box(
        slide,
        Inches(0.72),
        Inches(1.55),
        Inches(6.0),
        Inches(4.95),
        "프로젝트가 해결하려는 문제",
        [
            "치과 파노라마 X-ray 판독은 작은 병변, 중첩 구조, 해부학적 복잡성 때문에 해석 난도가 높다.",
            "환자 입장에서는 왜 특정 치료가 필요한지 시각적 근거를 이해하기 어렵고, 진료비 설명 부족에 대한 불신도 존재한다.",
            "임상의 판단을 대체하려는 것이 아니라, 병변 위치와 유형을 일관되게 제시하는 2차 의견 인터페이스가 필요하다.",
        ],
    )
    add_bullet_box(
        slide,
        Inches(6.93),
        Inches(1.55),
        Inches(5.65),
        Inches(4.95),
        "핵심 목표",
        [
            "파노라마 X-ray에서 병변 의심 부위를 bbox 기반으로 자동 검출한다.",
            "충치 계열 병변은 detector 내부 라벨과 무관하게 서비스 응답에서 `caries`로 통일한다.",
            "웹 화면에서 결과를 시각화하고, 향후 치료 옵션과 치료비 추정으로 확장 가능한 구조를 만든다.",
            "로컬 환경에서 학습, 평가, 서빙이 모두 가능한 local-first 워크플로우를 유지한다.",
        ],
    )
    add_footer(slide, 2)

    # 3. Strategy
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "APPROACH")
    add_title(slide, "접근 전략과 설계 원칙", "왜 2-stage 구조를 선택했는가")
    add_bullet_box(
        slide,
        Inches(0.72),
        Inches(1.48),
        Inches(4.0),
        Inches(4.95),
        "Detection 우선",
        [
            "분류(classification)만으로는 병변 위치를 설명할 수 없기 때문에 먼저 detection을 채택했다.",
            "bbox 위치를 제공하면 임상의와 환자 모두 어떤 부위를 근거로 판단했는지 확인할 수 있다.",
        ],
    )
    add_bullet_box(
        slide,
        Inches(4.87),
        Inches(1.48),
        Inches(4.0),
        Inches(4.95),
        "Hierarchical 설계",
        [
            "`deep_caries`는 표본 수가 너무 적어 서비스용 별도 클래스로 유지하기 어렵다.",
            "따라서 detector 내부 라벨이 `caries`/`deep_caries`/`caries_family`여도 서빙 응답은 `caries`로 정규화한다.",
        ],
    )
    add_bullet_box(
        slide,
        Inches(9.03),
        Inches(1.48),
        Inches(3.55),
        Inches(4.95),
        "설명 가능한 출력",
        [
            "웹 UI에서 병변 bbox, 클래스, 신뢰도, 부위별 치료비 추정 테이블을 함께 보여준다.",
            "최종 진단 확정이 아니라 의사결정 보조용 인터페이스임을 명확히 유지한다.",
        ],
    )
    add_footer(slide, 3)

    # 4. Architecture / sequential pipeline
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "ARCHITECTURE")
    add_title(slide, "전체 시스템 및 추론 흐름", "현재 구현 기준으로 detector와 classifier는 순차적으로 동작한다")
    add_pipeline_box(
        slide,
        Inches(1.65),
        [
            "웹 업로드",
            "Django /predict/",
            "YOLO Detector",
            "ROI Crop",
            "Response Normalize",
        ],
    )
    add_pipeline_box(
        slide,
        Inches(3.0),
        [
            "Result Merge",
            "JSON 응답 조립",
            "오버레이 시각화",
            "치료비 추정",
        ],
    )
    add_paragraph_box(
        slide,
        Inches(0.82),
        Inches(4.25),
        Inches(12.0),
        Inches(1.9),
        "핵심 설명",
        [
            "main detector는 `caries_family`, `periapical_lesion`, `impacted_tooth`, `retained_root`를 찾고, 서비스 응답에서는 충치 계열을 모두 `caries`로 통일한다.",
            "periodontal detector는 `bone_loss`, `furcation_involvement`를 별도로 찾고, 검출 ROI에 대해 후속 중증도 분류기를 적용한다.",
            "즉, 현재 서빙 구조는 충치 세분화가 아니라 서비스용 라벨 정규화와 치주 후속 분류를 결합한 형태다.",
        ],
    )
    add_footer(slide, 4)

    # 5. Datasets
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "DATASETS")
    add_title(slide, "사용 데이터셋 구성", "원본 형식과 프로젝트 내 역할")
    data = [
        ["데이터셋", "원본 형식", "프로젝트 역할", "주요 클래스 또는 매핑"],
        ["DENTEX", "JSON annotation", "기본 detection 라벨 소스\nseverity crop supervised source", "caries / deep_caries /\nperiapical_lesion / impacted_tooth"],
        ["CariesXrays", "Pascal VOC XML", "충치 계열 bbox 확장\npseudo-label 후보 crop source", "VOC `Decay` -> project `caries`"],
        ["UMFIH", "YOLO format", "병변 다양성 보강", "Carious lesion -> caries\nApical periodontitis -> periapical_lesion\nImpacted tooth -> impacted_tooth"],
    ]
    add_table_slide(
        slide,
        Inches(0.72),
        Inches(1.72),
        Inches(11.95),
        Inches(3.8),
        4,
        4,
        data,
        [Inches(1.6), Inches(1.65), Inches(3.4), Inches(5.3)],
    )
    add_paragraph_box(
        slide,
        Inches(0.72),
        Inches(5.78),
        Inches(11.95),
        Inches(0.95),
        "비고",
        [
            "현재 기본 detection 학습은 merged detection 데이터에서 파생된 hierarchical dataset을 사용하며, active split 규모는 train 5405 / val 745 / test 1316이다.",
        ],
    )
    add_footer(slide, 5)

    # 6. Preprocessing
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "PREPROCESSING")
    add_title(slide, "데이터 전처리 및 학습 파이프라인", "서로 다른 annotation 포맷을 하나의 workflow로 통합")
    add_bullet_box(
        slide,
        Inches(0.72),
        Inches(1.55),
        Inches(5.9),
        Inches(5.0),
        "Detection 전처리",
        [
            "DENTEX JSON을 YOLO bbox format으로 변환한다.",
            "CariesXrays VOC XML을 YOLO format으로 변환하고 `Decay`를 `caries`로 매핑한다.",
            "UMFIH YOLO 라벨 중 프로젝트와 안전하게 정렬되는 클래스만 remap하여 병합한다.",
            "필요 시 oversampling manifest를 만들어 희소 클래스 편향을 완화한다.",
        ],
    )
    add_bullet_box(
        slide,
        Inches(6.8),
        Inches(1.55),
        Inches(5.8),
        Inches(5.0),
        "Hierarchical + Severity 전처리",
        [
            "4-class detection 라벨을 3-class coarse 체계로 재구성하여 `caries_family`를 만든다.",
            "DENTEX lesion annotation을 crop으로 잘라 severity classifier용 분류 데이터셋을 구축한다.",
            "CariesXrays lesion crop은 unlabeled pool로 저장한 뒤 pseudo-labeling 실험에 활용할 수 있다.",
            "실행 스크립트는 `prepare_detection_dataset.py`, `prepare_cariesxrays_yolo.py`, `prepare_umfih_yolo.py`, `prepare_hierarchical_detection_dataset.py`, `prepare_severity_dataset.py` 순으로 구성된다.",
        ],
    )
    add_footer(slide, 6)

    # 7. Models
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "MODELS")
    add_title(slide, "모델 구성 및 주요 설정", "현재 기본값과 선정 이유")
    model_data = [
        ["구성 요소", "모델 / 설정", "선정 이유"],
        ["Detection", "Ultralytics YOLOv8n\nimgsz=416, batch=16, epochs=50,\nworkers=4, patience=10", "경량 모델로 로컬 GPU에서도 실험 가능하고, bbox 검출 속도와 정확도의 균형이 좋다."],
        ["Follow-up", "EfficientNet-B0 classifiers\ninput=224\nperiodontal severity stages", "치주 ROI 후속 분류 문제에 적합한 경량 CNN이며 서비스 응답과 수가 라우팅에 직접 연결된다."],
        ["Serving", "Django + DRF `/predict/`\n라벨 정규화 및 JSON 응답 조립", "웹 업로드, 추론, 시각화, 향후 진료비 추정 모듈까지 하나의 서비스 흐름으로 연결하기 쉽다."],
    ]
    add_table_slide(
        slide,
        Inches(0.72),
        Inches(1.68),
        Inches(11.95),
        Inches(4.3),
        4,
        3,
        model_data,
        [Inches(1.6), Inches(3.6), Inches(6.75)],
    )
    add_footer(slide, 7)

    # 8. Implementation / web
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "IMPLEMENTATION")
    add_title(slide, "웹 서비스 구현과 현재 동작", "사용자가 확인 가능한 산출물 중심")
    add_bullet_box(
        slide,
        Inches(0.72),
        Inches(1.55),
        Inches(5.8),
        Inches(4.9),
        "백엔드 및 API",
        [
            "`python manage.py runserver`로 Django 서비스를 실행한다.",
            "`GET /health/`로 상태 확인, `POST /predict/`로 multipart 이미지 업로드 추론을 처리한다.",
            "응답에는 detection bbox, class_name, confidence, 치주 severity 관련 메타데이터가 포함된다.",
        ],
    )
    add_bullet_box(
        slide,
        Inches(6.82),
        Inches(1.55),
        Inches(5.75),
        Inches(4.9),
        "프런트엔드 UI",
        [
            "업로드 이미지 미리보기와 오버레이 bbox를 한 화면에서 확인할 수 있다.",
            "탐지 목록에 병명, 부위, 신뢰도를 제시하고, 치료비 추정 표에 부위별 병명과 예상 비용을 정리한다.",
            "현재 치료비는 시연용 규칙 기반 값이며, 향후 전문가 자문 기반 fee table로 고도화할 계획이다.",
        ],
    )
    add_footer(slide, 8)

    # 9. Evaluation
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "EVALUATION")
    add_title(slide, "성능 평가 지표와 실험 계획", "실험 재현성과 비교 원칙")
    add_bullet_box(
        slide,
        Inches(0.72),
        Inches(1.58),
        Inches(3.9),
        Inches(4.8),
        "Detection 지표",
        [
            "Precision",
            "Recall",
            "mAP50",
            "mAP50-95",
            "핵심 기준은 bbox 위치 정확도까지 반영하는 `mAP50-95`다.",
        ],
    )
    add_bullet_box(
        slide,
        Inches(4.82),
        Inches(1.58),
        Inches(3.75),
        Inches(4.8),
        "Severity 지표",
        [
            "val_loss",
            "accuracy",
            "macro F1",
            "best checkpoint 선택과 early stopping은 현재 `val_loss` 기준이다.",
        ],
    )
    add_bullet_box(
        slide,
        Inches(8.78),
        Inches(1.58),
        Inches(3.8),
        Inches(4.8),
        "비교 실험 축",
        [
            "flat 4-class vs hierarchical 3-class detection",
            "imgsz=416 vs 512",
            "class balancing on/off",
            "periodontal follow-up classifier 사용 여부",
        ],
    )
    add_footer(slide, 9)

    # 10. Current status
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "STATUS")
    add_title(slide, "현재 구현 상태와 산출물", "문서, 코드, 모델, UI 기준 정리")
    add_paragraph_box(
        slide,
        Inches(0.72),
        Inches(1.52),
        Inches(12.0),
        Inches(5.05),
        "현재 기준 산출물",
        [
            "데이터 준비 스크립트: DENTEX, CariesXrays, UMFIH를 YOLO 기반 detection 파이프라인으로 정리하는 스크립트가 준비되어 있다.",
            "학습 파이프라인: hierarchical detection dataset 자동 준비, YOLOv8n 학습, periodontal follow-up classifier 학습 흐름이 README에 정리되어 있다.",
            "서빙 파이프라인: Django `/predict/` API가 detector 결과를 받은 뒤 충치 계열 응답을 `caries`로 정규화하고 치주 ROI에만 후속 분류를 수행한다.",
            "웹 인터페이스: 업로드, 탐지 결과 시각화, 탐지 목록, 치료비 추정 테이블, 공지 문구 등 제출 시연용 화면이 구현되어 있다.",
            "문서화: README와 ARCHITECTURE 문서에 데이터셋, 모델, 평가, 개발 일정, 참고문헌이 정리되어 있다.",
        ],
    )
    add_footer(slide, 10)

    # 11. Limitations
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "LIMITATIONS")
    add_title(slide, "한계점과 향후 확장 계획", "제출 시 명확히 밝혀야 할 부분")
    add_bullet_box(
        slide,
        Inches(0.72),
        Inches(1.55),
        Inches(5.9),
        Inches(4.95),
        "현재 한계",
        [
            "웹 서비스는 진단 확정 도구가 아니라 진단 보조용 인터페이스다.",
            "치료비 추정은 아직 규칙 기반 시범 기능이며, 실제 진료비는 병원과 처치 방식에 따라 달라질 수 있다.",
            "치주 follow-up classifier와 detector 모두 추가 임상 검증과 외부 데이터 평가가 필요하다.",
        ],
    )
    add_bullet_box(
        slide,
        Inches(6.82),
        Inches(1.55),
        Inches(5.75),
        Inches(4.95),
        "향후 확장",
        [
            "치과의사 자문 기반 진료비 계산 모듈 정교화",
            "환자 단위 리포트 자동 생성",
            "active learning 기반 재라벨링 루프",
            "다기관 데이터 추가와 일반화 성능 검증",
            "웹 대시보드와 실제 진료 워크플로우 연계",
        ],
    )
    add_footer(slide, 11)

    # 12. References
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_header_bar(slide, "REFERENCES")
    add_title(slide, "참고문헌", "제출 자료에 포함한 핵심 출처")
    add_paragraph_box(
        slide,
        Inches(0.72),
        Inches(1.45),
        Inches(12.0),
        Inches(5.25),
        "주요 참고 자료",
        [
            "Hamamci et al. (2023). DENTEX: An abnormal tooth detection benchmark for panoramic X-rays. arXiv.",
            "Chen et al. (2024). CariesXrays: Enhancing caries detection in hospital-scale panoramic dental X-rays via feature pyramid contrastive learning. AAAI.",
            "Mureșanu et al. (2025). Dataset for automating dental condition detection on panoramic radiographs. Zenodo.",
            "Tan & Le (2019). EfficientNet: Rethinking model scaling for convolutional neural networks. ICML.",
            "Ultralytics YOLO documentation, Django documentation, Django REST framework documentation.",
            "Korea Consumer Agency (2025). Dental treatment cost disputes are rapidly increasing; provision of treatment cost plans needs to be activated.",
        ],
    )
    add_footer(slide, 12)

    return prs


def main() -> None:
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_presentation()
    prs.save(OUT_PATH)
    print(OUT_PATH.resolve())


if __name__ == "__main__":
    main()
